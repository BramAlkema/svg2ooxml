from __future__ import annotations

from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402
from pptx.enum.dml import MSO_FILL  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402

from tests.visual.helpers.builder import PptxBuilder

FIXTURE = Path("tests/visual/fixtures/simple_rect.svg")


@pytest.mark.integration
def test_simple_rect_text_alignment(tmp_path) -> None:
    builder = PptxBuilder()
    svg_text = FIXTURE.read_text(encoding="utf-8")

    pptx_path = tmp_path / "simple_rect.pptx"
    build_result = builder.build_from_svg(svg_text, pptx_path)
    assert build_result.slide_count == 1

    presentation = Presentation(build_result.pptx_path)
    slide = presentation.slides[0]

    text_shape = None
    blue_rect = None
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if shape.has_text_frame and "svg2ooxml" in shape.text_frame.text:
            text_shape = shape
        if (
            shape.fill.type == MSO_FILL.SOLID
            and str(shape.fill.fore_color.rgb).upper() == "4285F4"
        ):
            blue_rect = shape

    assert text_shape is not None, "Expected a text shape containing 'svg2ooxml'."
    paragraph = text_shape.text_frame.paragraphs[0]
    assert paragraph.alignment == PP_ALIGN.CENTER, "Card headline should be centred."

    assert blue_rect is not None, "Expected primary blue rectangle."
    if blue_rect.line and blue_rect.line.fill.type == MSO_FILL.SOLID:
        assert str(blue_rect.line.fill.fore_color.rgb).upper() == "0F3AA3"
