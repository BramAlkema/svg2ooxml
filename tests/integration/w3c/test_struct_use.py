from __future__ import annotations

from pathlib import Path

import pytest
pptx = pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402
from pptx.enum.dml import MSO_FILL  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

from tools.visual.builder import PptxBuilder

W3C_FIXTURE = Path("tests/svg/struct-use-10-f.svg")


@pytest.mark.integration
def test_struct_use_rectangles_are_green(tmp_path) -> None:
    builder = PptxBuilder()
    svg_text = W3C_FIXTURE.read_text(encoding="utf-8")
    pptx_path = tmp_path / "struct_use.pptx"
    build_result = builder.build_from_svg(svg_text, pptx_path)

    assert build_result.slide_count == 1
    prs = Presentation(build_result.pptx_path)
    slide = prs.slides[0]

    greens: list[str] = []
    strokes: list[str] = []

    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if shape.has_text_frame and shape.text_frame and shape.text_frame.text.strip():
            continue

        if shape.fill.type == MSO_FILL.SOLID:
            greens.append(str(shape.fill.fore_color.rgb))
        if shape.line and shape.line.fill.type == MSO_FILL.SOLID:
            strokes.append(str(shape.line.fill.fore_color.rgb))

    assert greens.count("008000") >= 3, f"Expected at least three green rectangles, found {greens}"
    assert strokes.count("006400") >= 1, f"Expected a dark-green stroke in the collection, found {strokes}"
