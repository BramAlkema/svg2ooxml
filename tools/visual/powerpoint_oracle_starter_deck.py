#!/usr/bin/env python3
"""Build an editable PowerPoint oracle starter deck for groups and triggers.

The deck is intentionally small and opinionated:

- real grouped shapes, not fake layout clusters
- trigger scenarios authored as native PresentationML timing
- patterns borrowed from the mined oracle decks in docs/research/

The output is meant to be opened in PowerPoint and then refined manually.
"""

from __future__ import annotations

import argparse
import logging
import shutil
import tempfile
import zipfile
from collections.abc import Iterable, Iterator
from dataclasses import dataclass
from itertools import count
from pathlib import Path

from lxml import etree as ET

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - dependency check
    Presentation = None  # type: ignore[assignment]
    RGBColor = None  # type: ignore[assignment]
    MSO_AUTO_SHAPE_TYPE = None  # type: ignore[assignment]
    PP_ALIGN = None  # type: ignore[assignment]
    Inches = None  # type: ignore[assignment]
    Pt = None  # type: ignore[assignment]
    _PPTX_IMPORT_ERROR = exc
else:
    _PPTX_IMPORT_ERROR = None

from svg2ooxml.drawingml.xml_builder import NS_P, p_elem, p_sub

logger = logging.getLogger(__name__)

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5


@dataclass(frozen=True, slots=True)
class StartCondition:
    delay_ms: int = 0
    event: str | None = None
    target_shape_id: int | None = None


@dataclass(frozen=True, slots=True)
class BuildEntry:
    shape_id: int
    grp_id: int | str


@dataclass(frozen=True, slots=True)
class InteractiveSequence:
    trigger_shape_ids: tuple[int, ...]
    effect_pars: tuple[ET._Element, ...]


@dataclass(frozen=True, slots=True)
class SlideArtifact:
    timing_xml: str | None


def _require_python_pptx() -> None:
    if Presentation is None or Inches is None or Pt is None:
        raise RuntimeError(
            "python-pptx is required to build the PowerPoint oracle starter deck."
        ) from _PPTX_IMPORT_ERROR


def _add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int,
    bold: bool = False,
    color: tuple[int, int, int] = (0, 0, 0),
    align=None,
):
    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def _style_shape(
    shape,
    *,
    fill_rgb: tuple[int, int, int],
    line_rgb: tuple[int, int, int] = (45, 45, 45),
    line_width_pt: float = 1.25,
    text_rgb: tuple[int, int, int] = (0, 0, 0),
    font_size: int = 15,
    bold: bool = False,
) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(*line_rgb)
    shape.line.width = Pt(line_width_pt)
    if not shape.has_text_frame:
        return
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = shape.name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*text_rgb)


def _add_named_box(
    shape_container,
    *,
    name: str,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_rgb: tuple[int, int, int],
    line_rgb: tuple[int, int, int] = (45, 45, 45),
    font_size: int = 15,
    bold: bool = False,
):
    shape = shape_container.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.name = name
    _style_shape(
        shape,
        fill_rgb=fill_rgb,
        line_rgb=line_rgb,
        text_rgb=(0, 0, 0),
        font_size=font_size,
        bold=bold,
    )
    return shape


def _add_named_circle(
    shape_container,
    *,
    name: str,
    left: float,
    top: float,
    diameter: float,
    fill_rgb: tuple[int, int, int],
):
    shape = shape_container.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(left),
        Inches(top),
        Inches(diameter),
        Inches(diameter),
    )
    shape.name = name
    _style_shape(shape, fill_rgb=fill_rgb, font_size=14, bold=True)
    return shape


def _add_oracle_bar(
    slide,
    *,
    name: str,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_rgb: tuple[int, int, int],
    line_rgb: tuple[int, int, int],
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(*line_rgb)
    shape.line.width = Pt(1.5)
    return shape


def _add_oracle_marker(
    slide,
    *,
    left: float,
    top: float,
    height: float,
    label: str,
) -> None:
    marker = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.025),
        Inches(height),
    )
    marker.name = f"Marker {label}"
    marker.fill.solid()
    marker.fill.fore_color.rgb = RGBColor(185, 28, 28)
    marker.line.color.rgb = RGBColor(185, 28, 28)
    _add_textbox(
        slide,
        left - 0.18,
        top + height + 0.08,
        0.55,
        0.25,
        label,
        font_size=9,
        color=(31, 41, 55),
        align=PP_ALIGN.CENTER,
    )


def _add_group_triplet(
    slide,
    *,
    left: float,
    top: float,
    prefix: str,
):
    group = slide.shapes.add_group_shape()
    a = _add_named_box(
        group,
        name=f"{prefix} A",
        left=left,
        top=top,
        width=1.35,
        height=0.8,
        fill_rgb=(245, 208, 95),
        bold=True,
    )
    b = _add_named_box(
        group,
        name=f"{prefix} B",
        left=left + 1.6,
        top=top,
        width=1.35,
        height=0.8,
        fill_rgb=(147, 197, 253),
        bold=True,
    )
    c = _add_named_box(
        group,
        name=f"{prefix} C",
        left=left + 0.8,
        top=top + 1.05,
        width=1.35,
        height=0.8,
        fill_rgb=(167, 243, 208),
        bold=True,
    )
    group.name = f"{prefix} Group"
    return group, (a, b, c)


def _append_conditions(parent: ET._Element, conditions: Iterable[StartCondition]) -> None:
    for condition in conditions:
        attrs: dict[str, str] = {"delay": str(condition.delay_ms)}
        if condition.event:
            attrs["evt"] = condition.event
        cond = p_sub(parent, "cond", **attrs)
        if condition.target_shape_id is not None:
            tgt_el = p_sub(cond, "tgtEl")
            p_sub(tgt_el, "spTgt", spid=str(condition.target_shape_id))


def _build_set_visibility(
    *,
    id_counter: Iterator[int],
    target_shape_id: int,
    visibility: str,
) -> ET._Element:
    set_elem = p_elem("set")
    c_bhvr = p_sub(set_elem, "cBhvr")
    c_tn = p_sub(
        c_bhvr,
        "cTn",
        id=str(next(id_counter)),
        dur="1",
        fill="hold",
    )
    st_cond_lst = p_sub(c_tn, "stCondLst")
    p_sub(st_cond_lst, "cond", delay="0")
    tgt_el = p_sub(c_bhvr, "tgtEl")
    p_sub(tgt_el, "spTgt", spid=str(target_shape_id))
    attr_name_lst = p_sub(c_bhvr, "attrNameLst")
    p_sub(attr_name_lst, "attrName").text = "style.visibility"
    to = p_sub(set_elem, "to")
    p_sub(to, "strVal", val=visibility)
    return set_elem


def _build_anim_effect(
    *,
    id_counter: Iterator[int],
    target_shape_id: int,
    duration_ms: int,
    transition: str,
    filter_name: str,
) -> ET._Element:
    anim_effect = p_elem(
        "animEffect",
        transition=transition,
        filter=filter_name,
    )
    c_bhvr = p_sub(anim_effect, "cBhvr")
    p_sub(c_bhvr, "cTn", id=str(next(id_counter)), dur=str(duration_ms))
    tgt_el = p_sub(c_bhvr, "tgtEl")
    p_sub(tgt_el, "spTgt", spid=str(target_shape_id))
    return anim_effect


def _build_anim_scale(
    *,
    id_counter: Iterator[int],
    target_shape_id: int,
    duration_ms: int,
    by_x: int,
    by_y: int,
) -> ET._Element:
    anim_scale = p_elem("animScale")
    c_bhvr = p_sub(anim_scale, "cBhvr")
    p_sub(
        c_bhvr,
        "cTn",
        id=str(next(id_counter)),
        dur=str(duration_ms),
        fill="hold",
    )
    tgt_el = p_sub(c_bhvr, "tgtEl")
    p_sub(tgt_el, "spTgt", spid=str(target_shape_id))
    p_sub(anim_scale, "by", x=str(by_x), y=str(by_y))
    return anim_scale


def _build_native_width_scale_segment(
    *,
    id_counter: Iterator[int],
    target_shape_id: int,
    duration_ms: int,
    delay_ms: int,
    grp_id: int,
    by_x: int,
    by_y: int = 100000,
) -> ET._Element:
    return _build_effect_par(
        id_counter=id_counter,
        duration_ms=duration_ms,
        delay_ms=delay_ms,
        node_type="withEffect",
        preset_id=6,
        preset_class="emph",
        preset_subtype=0,
        grp_id=grp_id,
        child_elements=[
            _build_anim_scale(
                id_counter=id_counter,
                target_shape_id=target_shape_id,
                duration_ms=duration_ms,
                by_x=by_x,
                by_y=by_y,
            )
        ],
    )


def _build_anim_motion(
    *,
    id_counter: Iterator[int],
    target_shape_id: int,
    duration_ms: int,
    path: str,
) -> ET._Element:
    anim_motion = p_elem(
        "animMotion",
        origin="layout",
        path=path,
        pathEditMode="relative",
        ptsTypes="",
    )
    c_bhvr = p_sub(anim_motion, "cBhvr")
    p_sub(
        c_bhvr,
        "cTn",
        id=str(next(id_counter)),
        dur=str(duration_ms),
        fill="hold",
    )
    tgt_el = p_sub(c_bhvr, "tgtEl")
    p_sub(tgt_el, "spTgt", spid=str(target_shape_id))
    attr_name_lst = p_sub(c_bhvr, "attrNameLst")
    p_sub(attr_name_lst, "attrName").text = "ppt_x"
    p_sub(attr_name_lst, "attrName").text = "ppt_y"
    return anim_motion


def _build_effect_par(
    *,
    id_counter: Iterator[int],
    duration_ms: int,
    node_type: str,
    child_elements: Iterable[ET._Element],
    delay_ms: int = 0,
    preset_id: int | None = None,
    preset_class: str | None = None,
    preset_subtype: int | None = 0,
    grp_id: int | str = "0",
    start_conditions: Iterable[StartCondition] | None = None,
) -> ET._Element:
    par = p_elem("par")
    c_tn_attrs: dict[str, str] = {
        "id": str(next(id_counter)),
        "dur": str(duration_ms),
        "fill": "hold",
        "grpId": grp_id,
        "nodeType": node_type,
    }
    if preset_id is not None:
        c_tn_attrs["presetID"] = str(preset_id)
    if preset_class is not None:
        c_tn_attrs["presetClass"] = preset_class
    if preset_subtype is not None and preset_id is not None:
        c_tn_attrs["presetSubtype"] = str(preset_subtype)
    c_tn = p_sub(par, "cTn", **c_tn_attrs)
    st_cond_lst = p_sub(c_tn, "stCondLst")
    conditions = tuple(start_conditions or (StartCondition(delay_ms=delay_ms),))
    _append_conditions(st_cond_lst, conditions)
    child_tn_lst = p_sub(c_tn, "childTnLst")
    for child in child_elements:
        child_tn_lst.append(child)
    return par


def _build_interactive_seq(
    *,
    id_counter: Iterator[int],
    trigger_shape_ids: Iterable[int],
    effect_pars: Iterable[ET._Element],
) -> ET._Element:
    trigger_ids = tuple(trigger_shape_ids)
    seq = p_elem("seq", concurrent="1", nextAc="seek")
    c_tn = p_sub(
        seq,
        "cTn",
        id=str(next(id_counter)),
        restart="whenNotActive",
        fill="hold",
        evtFilter="cancelBubble",
        nodeType="interactiveSeq",
    )
    st_cond_lst = p_sub(c_tn, "stCondLst")
    _append_conditions(
        st_cond_lst,
        [
            StartCondition(event="onClick", delay_ms=0, target_shape_id=shape_id)
            for shape_id in trigger_ids
        ],
    )
    end_sync = p_sub(c_tn, "endSync", evt="end", delay="0")
    p_sub(end_sync, "rtn", val="all")
    child_tn_lst = p_sub(c_tn, "childTnLst")
    outer_par = p_sub(child_tn_lst, "par")
    outer_ctn = p_sub(outer_par, "cTn", id=str(next(id_counter)), fill="hold")
    outer_st = p_sub(outer_ctn, "stCondLst")
    p_sub(outer_st, "cond", delay="0")
    outer_children = p_sub(outer_ctn, "childTnLst")
    inner_par = p_sub(outer_children, "par")
    inner_ctn = p_sub(inner_par, "cTn", id=str(next(id_counter)), fill="hold")
    inner_st = p_sub(inner_ctn, "stCondLst")
    p_sub(inner_st, "cond", delay="0")
    inner_children = p_sub(inner_ctn, "childTnLst")
    for effect_par in effect_pars:
        inner_children.append(effect_par)
    next_cond_lst = p_sub(seq, "nextCondLst")
    _append_conditions(
        next_cond_lst,
        [
            StartCondition(event="onClick", delay_ms=0, target_shape_id=shape_id)
            for shape_id in trigger_ids
        ],
    )
    return seq


def _build_timing_xml(
    *,
    start_id: int,
    main_effect_pars: Iterable[ET._Element] = (),
    interactive_sequences: Iterable[InteractiveSequence] = (),
    build_entries: Iterable[BuildEntry] = (),
) -> str:
    id_counter = count(start_id)
    timing = p_elem("timing")
    tn_lst = p_sub(timing, "tnLst")
    root_par = p_sub(tn_lst, "par")
    root_ctn = p_sub(
        root_par,
        "cTn",
        id=str(next(id_counter)),
        dur="indefinite",
        restart="never",
        nodeType="tmRoot",
    )
    root_children = p_sub(root_ctn, "childTnLst")

    effect_pars = tuple(main_effect_pars)
    if effect_pars:
        seq = p_sub(root_children, "seq", concurrent="1", nextAc="seek")
        seq_ctn = p_sub(
            seq,
            "cTn",
            id=str(next(id_counter)),
            dur="indefinite",
            nodeType="mainSeq",
        )
        seq_children = p_sub(seq_ctn, "childTnLst")
        outer_par = p_sub(seq_children, "par")
        outer_ctn = p_sub(outer_par, "cTn", id=str(next(id_counter)), fill="hold")
        outer_st = p_sub(outer_ctn, "stCondLst")
        p_sub(outer_st, "cond", delay="indefinite")
        outer_children = p_sub(outer_ctn, "childTnLst")
        inner_par = p_sub(outer_children, "par")
        inner_ctn = p_sub(inner_par, "cTn", id=str(next(id_counter)), fill="hold")
        inner_st = p_sub(inner_ctn, "stCondLst")
        p_sub(inner_st, "cond", delay="0")
        inner_children = p_sub(inner_ctn, "childTnLst")
        for effect_par in effect_pars:
            inner_children.append(effect_par)
        prev_cond_lst = p_sub(seq, "prevCondLst")
        prev_cond = p_sub(prev_cond_lst, "cond", evt="onPrev", delay="0")
        p_sub(p_sub(prev_cond, "tgtEl"), "sldTgt")
        next_cond_lst = p_sub(seq, "nextCondLst")
        next_cond = p_sub(next_cond_lst, "cond", evt="onNext", delay="0")
        p_sub(p_sub(next_cond, "tgtEl"), "sldTgt")

    for interactive in interactive_sequences:
        root_children.append(
            _build_interactive_seq(
                id_counter=id_counter,
                trigger_shape_ids=interactive.trigger_shape_ids,
                effect_pars=interactive.effect_pars,
            )
        )

    entries = tuple(build_entries)
    if entries:
        bld_lst = p_sub(timing, "bldLst")
        for entry in entries:
            p_sub(
                bld_lst,
                "bldP",
                spid=str(entry.shape_id),
                grpId=entry.grp_id,
                animBg="1",
            )

    return ET.tostring(timing, encoding="unicode")


def _patch_slide_xml_with_timing(slide_xml: bytes, timing_xml: str) -> bytes:
    slide_root = ET.fromstring(slide_xml)
    for existing in slide_root.findall(f"{{{NS_P}}}timing"):
        slide_root.remove(existing)
    slide_root.append(ET.fromstring(timing_xml.encode("utf-8")))
    return ET.tostring(
        slide_root,
        encoding="utf-8",
        xml_declaration=True,
        standalone="yes",
    )


def _inject_timing_map_into_pptx(pptx_path: Path, timing_by_slide_number: dict[int, str]) -> None:
    temp_dir = Path(tempfile.mkdtemp(prefix="ppt-oracle-starter-"))
    temp_pptx = temp_dir / pptx_path.name
    try:
        with zipfile.ZipFile(pptx_path, "r") as source, zipfile.ZipFile(
            temp_pptx,
            "w",
            compression=zipfile.ZIP_DEFLATED,
        ) as target:
            for info in source.infolist():
                payload = source.read(info.filename)
                if info.filename.startswith("ppt/slides/slide") and info.filename.endswith(".xml"):
                    slide_name = Path(info.filename).stem
                    slide_number = int(slide_name.replace("slide", ""))
                    timing_xml = timing_by_slide_number.get(slide_number)
                    if timing_xml:
                        payload = _patch_slide_xml_with_timing(payload, timing_xml)
                target.writestr(info, payload)
        shutil.move(temp_pptx, pptx_path)
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def _oracle_note(slide, source: str, summary: str) -> None:
    _add_textbox(
        slide,
        0.55,
        0.25,
        12.2,
        0.45,
        source,
        font_size=13,
        bold=True,
        color=(55, 65, 81),
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        0.55,
        0.62,
        12.2,
        0.42,
        summary,
        font_size=11,
        color=(75, 85, 99),
        align=PP_ALIGN.CENTER,
    )


def _title_slide(presentation) -> SlideArtifact:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_textbox(
        slide,
        0.6,
        0.55,
        12.1,
        0.55,
        "PowerPoint Oracle Starter: Groups and Triggers",
        font_size=26,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        1.0,
        1.3,
        11.3,
        0.55,
        "Built from mined PowerPoint timing patterns. Open this in PowerPoint, tweak, then feed the edited deck back into the oracle extractor.",
        font_size=14,
        color=(55, 65, 81),
        align=PP_ALIGN.CENTER,
    )
    bullet_y = 2.15
    bullets = [
        "Slides 2-5: group visibility and child animation inside a real group.",
        "Slides 6-8: click triggers authored as interactive sequences.",
        "Slides 9-10: event chaining probes for onBegin and onEnd.",
        "Slide 11: one target effect listening to two trigger shapes.",
        "Slide 12: native width scale semantics learned from W3C calcMode cases.",
    ]
    for bullet in bullets:
        _add_textbox(
            slide,
            1.2,
            bullet_y,
            10.9,
            0.35,
            f"- {bullet}",
            font_size=15,
            color=(31, 41, 55),
        )
        bullet_y += 0.45
    return SlideArtifact(timing_xml=None)


def _group_reveal_slide(presentation) -> SlideArtifact:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _oracle_note(
        slide,
        "Oracle base: example2 slide2, simplified to a single group entrance.",
        "Click once in slideshow: the whole group should appear together.",
    )
    group, _ = _add_group_triplet(slide, left=4.55, top=2.15, prefix="Reveal")
    max_shape_id = max(shape.shape_id for shape in slide.shapes)
    id_counter = count(max_shape_id + 5)
    grp_id = 1
    effect = _build_effect_par(
        id_counter=id_counter,
        duration_ms=1,
        node_type="clickEffect",
        preset_id=1,
        preset_class="entr",
        grp_id=grp_id,
        child_elements=[
            _build_set_visibility(
                id_counter=id_counter,
                target_shape_id=group.shape_id,
                visibility="visible",
            )
        ],
    )
    timing_xml = _build_timing_xml(
        start_id=max_shape_id + 1,
        main_effect_pars=[effect],
        build_entries=[BuildEntry(shape_id=group.shape_id, grp_id=grp_id)],
    )
    return SlideArtifact(timing_xml=timing_xml)


def _group_hide_slide(presentation) -> SlideArtifact:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _oracle_note(
        slide,
        "Oracle base: example2 slide2, isolated exit case.",
        "Click once in slideshow: the whole group should disappear together.",
    )
    group, _ = _add_group_triplet(slide, left=4.55, top=2.15, prefix="Hide")
    max_shape_id = max(shape.shape_id for shape in slide.shapes)
    id_counter = count(max_shape_id + 5)
    grp_id = 1
    effect = _build_effect_par(
        id_counter=id_counter,
        duration_ms=1,
        node_type="clickEffect",
        preset_id=1,
        preset_class="exit",
        grp_id=grp_id,
        child_elements=[
            _build_set_visibility(
                id_counter=id_counter,
                target_shape_id=group.shape_id,
                visibility="hidden",
            )
        ],
    )
    timing_xml = _build_timing_xml(
        start_id=max_shape_id + 1,
        main_effect_pars=[effect],
        build_entries=[BuildEntry(shape_id=group.shape_id, grp_id=grp_id)],
    )
    return SlideArtifact(timing_xml=timing_xml)


def _group_reveal_hide_reveal_slide(presentation) -> SlideArtifact:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _oracle_note(
        slide,
        "Oracle base: example2 slide2 appear -> disappear -> appear.",
        "One click should start all three phases: reveal, hide, reveal.",
    )
    group, _ = _add_group_triplet(slide, left=4.55, top=2.15, prefix="Cycle")
    max_shape_id = max(shape.shape_id for shape in slide.shapes)
    id_counter = count(max_shape_id + 5)
    grp_ids = (1, 2, 3)
    effects = [
        _build_effect_par(
            id_counter=id_counter,
            duration_ms=1,
            node_type="clickEffect",
            preset_id=1,
            preset_class="entr",
            grp_id=grp_ids[0],
            child_elements=[
                _build_set_visibility(
                    id_counter=id_counter,
                    target_shape_id=group.shape_id,
                    visibility="visible",
                )
            ],
        ),
        _build_effect_par(
            id_counter=id_counter,
            duration_ms=1,
            node_type="withEffect",
            delay_ms=500,
            preset_id=1,
            preset_class="exit",
            grp_id=grp_ids[1],
            child_elements=[
                _build_set_visibility(
                    id_counter=id_counter,
                    target_shape_id=group.shape_id,
                    visibility="hidden",
                )
            ],
        ),
        _build_effect_par(
            id_counter=id_counter,
            duration_ms=1,
            node_type="withEffect",
            delay_ms=500,
            preset_id=1,
            preset_class="entr",
            grp_id=grp_ids[2],
            child_elements=[
                _build_set_visibility(
                    id_counter=id_counter,
                    target_shape_id=group.shape_id,
                    visibility="visible",
                )
            ],
        ),
    ]
    timing_xml = _build_timing_xml(
        start_id=max_shape_id + 1,
        main_effect_pars=effects,
        build_entries=[
            BuildEntry(shape_id=group.shape_id, grp_id=grp_ids[0]),
            BuildEntry(shape_id=group.shape_id, grp_id=grp_ids[1]),
            BuildEntry(shape_id=group.shape_id, grp_id=grp_ids[2]),
        ],
    )
    return SlideArtifact(timing_xml=timing_xml)


def _group_reveal_plus_child_motion_slide(presentation) -> SlideArtifact:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _oracle_note(
        slide,
        "Oracle base: example2 slide2 plus example3 slide4 path effect.",
        "One click should reveal the group, then a child inside the group should move.",
    )
    group, children = _add_group_triplet(slide, left=4.25, top=2.0, prefix="Move")
    max_shape_id = max(shape.shape_id for shape in slide.shapes)
    id_counter = count(max_shape_id + 5)
    reveal = _build_effect_par(
        id_counter=id_counter,
        duration_ms=1,
        node_type="clickEffect",
        preset_id=1,
        preset_class="entr",
        grp_id=1,
        child_elements=[
            _build_set_visibility(
                id_counter=id_counter,
                target_shape_id=group.shape_id,
                visibility="visible",
            )
        ],
    )
    motion = _build_effect_par(
        id_counter=id_counter,
        duration_ms=1800,
        delay_ms=350,
        node_type="withEffect",
        preset_id=31,
        preset_class="path",
        grp_id=0,
        child_elements=[
            _build_anim_motion(
                id_counter=id_counter,
                target_shape_id=children[1].shape_id,
                duration_ms=1800,
                path="M 0 0 C 0.05 -0.1 0.16 -0.08 0.2 0.03 L 0.25 0.12 E",
            )
        ],
    )
    timing_xml = _build_timing_xml(
        start_id=max_shape_id + 1,
        main_effect_pars=[reveal, motion],
        build_entries=[
            BuildEntry(shape_id=group.shape_id, grp_id=1),
            BuildEntry(shape_id=children[1].shape_id, grp_id=0),
        ],
    )
    return SlideArtifact(timing_xml=timing_xml)


def _click_self_pulse_slide(presentation) -> SlideArtifact:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _oracle_note(
        slide,
        "Oracle base: example3 slide5 interactive self trigger.",
        "Click the blue button itself: it should pulse larger and return.",
    )
    button = _add_named_box(
        slide,
        name="Click Me",
        left=5.0,
        top=2.65,
        width=3.0,
        height=1.1,
        fill_rgb=(147, 197, 253),
        font_size=18,
        bold=True,
    )
    max_shape_id = max(shape.shape_id for shape in slide.shapes)
    id_counter = count(max_shape_id + 5)
    pulse = _build_effect_par(
        id_counter=id_counter,
        duration_ms=1800,
        node_type="clickEffect",
        preset_id=6,
        preset_class="emph",
        grp_id=0,
        child_elements=[
            _build_anim_scale(
                id_counter=id_counter,
                target_shape_id=button.shape_id,
                duration_ms=1800,
                by_x=180000,
                by_y=180000,
            )
        ],
    )
    timing_xml = _build_timing_xml(
        start_id=max_shape_id + 1,
        interactive_sequences=[
            InteractiveSequence(
                trigger_shape_ids=(button.shape_id,),
                effect_pars=(pulse,),
            )
        ],
        build_entries=[BuildEntry(shape_id=button.shape_id, grp_id=0)],
    )
    return SlideArtifact(timing_xml=timing_xml)


def _click_other_reveal_fade_slide(presentation) -> SlideArtifact:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _oracle_note(
        slide,
        "Oracle base: example3 slide5 click target -> set visibility + fade entrance.",
        "Click the orange trigger: the blue target should fade in.",
    )
    trigger = _add_named_box(
        slide,
        name="Trigger",
        left=1.75,
        top=2.8,
        width=2.25,
        height=0.95,
        fill_rgb=(253, 186, 116),
        font_size=16,
        bold=True,
    )
    target = _add_named_box(
        slide,
        name="Reveal Target",
        left=7.25,
        top=2.45,
        width=3.15,
        height=1.4,
        fill_rgb=(191, 219, 254),
        font_size=17,
        bold=True,
    )
    max_shape_id = max(shape.shape_id for shape in slide.shapes)
    id_counter = count(max_shape_id + 5)
    reveal = _build_effect_par(
        id_counter=id_counter,
        duration_ms=500,
        node_type="clickEffect",
        preset_id=9,
        preset_class="entr",
        grp_id=1,
        child_elements=[
            _build_set_visibility(
                id_counter=id_counter,
                target_shape_id=target.shape_id,
                visibility="visible",
            ),
            _build_anim_effect(
                id_counter=id_counter,
                target_shape_id=target.shape_id,
                duration_ms=500,
                transition="in",
                filter_name="fade",
            ),
        ],
    )
    timing_xml = _build_timing_xml(
        start_id=max_shape_id + 1,
        interactive_sequences=[
            InteractiveSequence(
                trigger_shape_ids=(trigger.shape_id,),
                effect_pars=(reveal,),
            )
        ],
        build_entries=[
            BuildEntry(shape_id=target.shape_id, grp_id=0),
            BuildEntry(shape_id=target.shape_id, grp_id=1),
        ],
    )
    return SlideArtifact(timing_xml=timing_xml)


def _click_other_motion_slide(presentation) -> SlideArtifact:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _oracle_note(
        slide,
        "Oracle base: example3 slide5 click target -> motion path.",
        "Click the orange trigger: the blue disc should follow the short path.",
    )
    trigger = _add_named_box(
        slide,
        name="Trigger Motion",
        left=1.35,
        top=2.85,
        width=2.7,
        height=0.95,
        fill_rgb=(253, 186, 116),
        font_size=15,
        bold=True,
    )
    dot = _add_named_circle(
        slide,
        name="Dot",
        left=7.05,
        top=2.8,
        diameter=0.95,
        fill_rgb=(96, 165, 250),
    )
    max_shape_id = max(shape.shape_id for shape in slide.shapes)
    id_counter = count(max_shape_id + 5)
    motion = _build_effect_par(
        id_counter=id_counter,
        duration_ms=1800,
        node_type="clickEffect",
        preset_id=31,
        preset_class="path",
        grp_id=0,
        child_elements=[
            _build_anim_motion(
                id_counter=id_counter,
                target_shape_id=dot.shape_id,
                duration_ms=1800,
                path="M 0 0 C 0.04 -0.12 0.17 -0.12 0.24 0.02 C 0.29 0.11 0.21 0.22 0.1 0.2 E",
            )
        ],
    )
    timing_xml = _build_timing_xml(
        start_id=max_shape_id + 1,
        interactive_sequences=[
            InteractiveSequence(
                trigger_shape_ids=(trigger.shape_id,),
                effect_pars=(motion,),
            )
        ],
        build_entries=[BuildEntry(shape_id=dot.shape_id, grp_id=0)],
    )
    return SlideArtifact(timing_xml=timing_xml)


def _on_other_begin_slide(presentation) -> SlideArtifact:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _oracle_note(
        slide,
        "Probe case: trigger event should be onBegin of another shape.",
        "Click once: the trigger box should pulse, and the blue target should reveal when that starts.",
    )
    trigger = _add_named_box(
        slide,
        name="Primary",
        left=2.0,
        top=2.75,
        width=2.4,
        height=0.95,
        fill_rgb=(253, 186, 116),
        font_size=16,
        bold=True,
    )
    target = _add_named_box(
        slide,
        name="On Begin",
        left=7.15,
        top=2.45,
        width=2.95,
        height=1.25,
        fill_rgb=(191, 219, 254),
        font_size=17,
        bold=True,
    )
    max_shape_id = max(shape.shape_id for shape in slide.shapes)
    id_counter = count(max_shape_id + 5)
    primary = _build_effect_par(
        id_counter=id_counter,
        duration_ms=1400,
        node_type="clickEffect",
        preset_id=6,
        preset_class="emph",
        grp_id=0,
        child_elements=[
            _build_anim_scale(
                id_counter=id_counter,
                target_shape_id=trigger.shape_id,
                duration_ms=1400,
                by_x=140000,
                by_y=140000,
            )
        ],
    )
    dependent = _build_effect_par(
        id_counter=id_counter,
        duration_ms=1,
        node_type="withEffect",
        preset_id=1,
        preset_class="entr",
        grp_id=1,
        start_conditions=[
            StartCondition(event="onBegin", delay_ms=0, target_shape_id=trigger.shape_id)
        ],
        child_elements=[
            _build_set_visibility(
                id_counter=id_counter,
                target_shape_id=target.shape_id,
                visibility="visible",
            )
        ],
    )
    timing_xml = _build_timing_xml(
        start_id=max_shape_id + 1,
        main_effect_pars=[primary, dependent],
        build_entries=[
            BuildEntry(shape_id=trigger.shape_id, grp_id=0),
            BuildEntry(shape_id=target.shape_id, grp_id=1),
        ],
    )
    return SlideArtifact(timing_xml=timing_xml)


def _on_other_end_slide(presentation) -> SlideArtifact:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _oracle_note(
        slide,
        "Probe case: trigger event should be onEnd of another shape.",
        "Click once: the trigger pulse should finish, then the blue target should reveal.",
    )
    trigger = _add_named_box(
        slide,
        name="Primary End",
        left=2.0,
        top=2.75,
        width=2.6,
        height=0.95,
        fill_rgb=(253, 186, 116),
        font_size=16,
        bold=True,
    )
    target = _add_named_box(
        slide,
        name="On End",
        left=7.15,
        top=2.45,
        width=2.8,
        height=1.25,
        fill_rgb=(191, 219, 254),
        font_size=17,
        bold=True,
    )
    max_shape_id = max(shape.shape_id for shape in slide.shapes)
    id_counter = count(max_shape_id + 5)
    primary = _build_effect_par(
        id_counter=id_counter,
        duration_ms=1400,
        node_type="clickEffect",
        preset_id=6,
        preset_class="emph",
        grp_id=0,
        child_elements=[
            _build_anim_scale(
                id_counter=id_counter,
                target_shape_id=trigger.shape_id,
                duration_ms=1400,
                by_x=140000,
                by_y=140000,
            )
        ],
    )
    dependent = _build_effect_par(
        id_counter=id_counter,
        duration_ms=1,
        node_type="withEffect",
        preset_id=1,
        preset_class="entr",
        grp_id=1,
        start_conditions=[
            StartCondition(event="onEnd", delay_ms=0, target_shape_id=trigger.shape_id)
        ],
        child_elements=[
            _build_set_visibility(
                id_counter=id_counter,
                target_shape_id=target.shape_id,
                visibility="visible",
            )
        ],
    )
    timing_xml = _build_timing_xml(
        start_id=max_shape_id + 1,
        main_effect_pars=[primary, dependent],
        build_entries=[
            BuildEntry(shape_id=trigger.shape_id, grp_id=0),
            BuildEntry(shape_id=target.shape_id, grp_id=1),
        ],
    )
    return SlideArtifact(timing_xml=timing_xml)


def _multi_trigger_slide(presentation) -> SlideArtifact:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _oracle_note(
        slide,
        "Probe case: one effect listens to two click triggers.",
        "Click either orange trigger: the blue target should pulse.",
    )
    trigger_a = _add_named_box(
        slide,
        name="Trigger A",
        left=1.55,
        top=2.65,
        width=2.15,
        height=0.9,
        fill_rgb=(253, 186, 116),
        font_size=15,
        bold=True,
    )
    trigger_b = _add_named_box(
        slide,
        name="Trigger B",
        left=1.55,
        top=3.75,
        width=2.15,
        height=0.9,
        fill_rgb=(251, 191, 36),
        font_size=15,
        bold=True,
    )
    target = _add_named_box(
        slide,
        name="Shared Target",
        left=7.0,
        top=3.0,
        width=3.2,
        height=1.2,
        fill_rgb=(191, 219, 254),
        font_size=17,
        bold=True,
    )
    max_shape_id = max(shape.shape_id for shape in slide.shapes)
    id_counter = count(max_shape_id + 5)
    pulse = _build_effect_par(
        id_counter=id_counter,
        duration_ms=1500,
        node_type="clickEffect",
        preset_id=6,
        preset_class="emph",
        grp_id=0,
        child_elements=[
            _build_anim_scale(
                id_counter=id_counter,
                target_shape_id=target.shape_id,
                duration_ms=1500,
                by_x=150000,
                by_y=150000,
            )
        ],
    )
    timing_xml = _build_timing_xml(
        start_id=max_shape_id + 1,
        interactive_sequences=[
            InteractiveSequence(
                trigger_shape_ids=(trigger_a.shape_id, trigger_b.shape_id),
                effect_pars=(pulse,),
            )
        ],
        build_entries=[BuildEntry(shape_id=target.shape_id, grp_id=0)],
    )
    return SlideArtifact(timing_xml=timing_xml)


def _native_width_scale_oracle_slide(presentation) -> SlideArtifact:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _oracle_note(
        slide,
        "Oracle addition: native width animScale semantics from W3C calcMode debugging.",
        "Click once: good factor shrink preserves height; delta-style shrink is the bad control; segmented shrink needs one build entry per segment.",
    )
    _add_textbox(
        slide,
        0.85,
        1.3,
        3.3,
        0.3,
        "Correct: by is a scale factor",
        font_size=12,
        bold=True,
        color=(31, 41, 55),
    )
    good = _add_oracle_bar(
        slide,
        name="Good by-factor shrink",
        left=0.9,
        top=1.72,
        width=3.2,
        height=0.55,
        fill_rgb=(68, 170, 255),
        line_rgb=(136, 0, 136),
    )
    _add_textbox(
        slide,
        0.9,
        2.38,
        3.9,
        0.34,
        '<p:by x="33333" y="100000"/>',
        font_size=10,
        color=(75, 85, 99),
    )

    _add_textbox(
        slide,
        4.95,
        1.3,
        3.5,
        0.3,
        "Bad control: by is not a delta",
        font_size=12,
        bold=True,
        color=(31, 41, 55),
    )
    bad = _add_oracle_bar(
        slide,
        name="Bad delta-style shrink",
        left=5.0,
        top=1.72,
        width=3.2,
        height=0.55,
        fill_rgb=(252, 165, 165),
        line_rgb=(153, 27, 27),
    )
    _add_textbox(
        slide,
        5.0,
        2.38,
        4.0,
        0.34,
        '<p:by x="-66667" y="0"/> collapses',
        font_size=10,
        color=(75, 85, 99),
    )

    _add_textbox(
        slide,
        0.85,
        3.25,
        7.0,
        0.3,
        "Segmented native scale: each segment has its own grpId + bldP",
        font_size=12,
        bold=True,
        color=(31, 41, 55),
    )
    segmented = _add_oracle_bar(
        slide,
        name="Segmented width shrink",
        left=0.9,
        top=3.76,
        width=4.8,
        height=0.65,
        fill_rgb=(68, 170, 255),
        line_rgb=(136, 0, 136),
    )
    marker_origin = 0.9
    marker_width = 4.8
    marker_top = 3.52
    marker_height = 1.08
    for value, label in ((300, "0s"), (255, "1.5s"), (180, "4s"), (30, "9s")):
        _add_oracle_marker(
            slide,
            left=marker_origin + marker_width * (value / 300.0),
            top=marker_top,
            height=marker_height,
            label=label,
        )
    _add_textbox(
        slide,
        0.9,
        4.88,
        8.6,
        0.45,
        "Segments: 300->255 uses x=85000, 255->180 uses x=70588, 180->30 uses x=16667. Untouched axis stays y=100000.",
        font_size=10,
        color=(75, 85, 99),
    )

    _add_textbox(
        slide,
        9.35,
        1.34,
        3.15,
        3.65,
        "What this oracle guards:\n"
        "- native only, no raster fallback\n"
        "- <p:by> is factor, not delta\n"
        "- width-only keeps y=100000\n"
        "- each segment needs animBg bldP\n"
        "- additive motion compensation is not used",
        font_size=11,
        color=(31, 41, 55),
    )

    max_shape_id = max(shape.shape_id for shape in slide.shapes)
    id_counter = count(max_shape_id + 5)
    good_grp_id = 1
    bad_grp_id = 2
    segment_grp_ids = (3, 4, 5)
    good_effect = _build_native_width_scale_segment(
        id_counter=id_counter,
        target_shape_id=good.shape_id,
        duration_ms=4000,
        delay_ms=0,
        grp_id=good_grp_id,
        by_x=33333,
    )
    bad_effect = _build_native_width_scale_segment(
        id_counter=id_counter,
        target_shape_id=bad.shape_id,
        duration_ms=4000,
        delay_ms=0,
        grp_id=bad_grp_id,
        by_x=-66667,
        by_y=0,
    )
    segmented_effects = [
        _build_native_width_scale_segment(
            id_counter=id_counter,
            target_shape_id=segmented.shape_id,
            duration_ms=1500,
            delay_ms=0,
            grp_id=segment_grp_ids[0],
            by_x=85000,
        ),
        _build_native_width_scale_segment(
            id_counter=id_counter,
            target_shape_id=segmented.shape_id,
            duration_ms=2500,
            delay_ms=1500,
            grp_id=segment_grp_ids[1],
            by_x=70588,
        ),
        _build_native_width_scale_segment(
            id_counter=id_counter,
            target_shape_id=segmented.shape_id,
            duration_ms=5000,
            delay_ms=4000,
            grp_id=segment_grp_ids[2],
            by_x=16667,
        ),
    ]
    timing_xml = _build_timing_xml(
        start_id=max_shape_id + 1,
        main_effect_pars=[good_effect, bad_effect, *segmented_effects],
        build_entries=[
            BuildEntry(shape_id=good.shape_id, grp_id=good_grp_id),
            BuildEntry(shape_id=bad.shape_id, grp_id=bad_grp_id),
            *(
                BuildEntry(shape_id=segmented.shape_id, grp_id=grp_id)
                for grp_id in segment_grp_ids
            ),
        ],
    )
    return SlideArtifact(timing_xml=timing_xml)


def build_oracle_starter_deck(output_path: Path) -> Path:
    _require_python_pptx()
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)

    slide_builders = [
        _title_slide,
        _group_reveal_slide,
        _group_hide_slide,
        _group_reveal_hide_reveal_slide,
        _group_reveal_plus_child_motion_slide,
        _click_self_pulse_slide,
        _click_other_reveal_fade_slide,
        _click_other_motion_slide,
        _on_other_begin_slide,
        _on_other_end_slide,
        _multi_trigger_slide,
        _native_width_scale_oracle_slide,
    ]
    artifacts = [builder(presentation) for builder in slide_builders]

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)

    timing_by_slide_number = {
        slide_number: artifact.timing_xml
        for slide_number, artifact in enumerate(artifacts, start=1)
        if artifact.timing_xml
    }
    _inject_timing_map_into_pptx(output_path, timing_by_slide_number)
    return output_path


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=Path("tmp/powerpoint-oracle-groups-triggers-starter.pptx"),
        help="Output PPTX path.",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Enable debug logging.",
    )
    return parser.parse_args()


def main() -> None:
    args = _parse_args()
    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(levelname)s %(message)s",
    )
    output = build_oracle_starter_deck(args.output)
    logger.info("Built PowerPoint oracle starter deck: %s", output)


if __name__ == "__main__":
    main()
