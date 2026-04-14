"""Shared helpers for hand-built animation sample PPTXs.

These helpers drive python-pptx to emit a minimal slide, then patch the
serialized ``slide1.xml`` with a hand-built ``<p:timing>`` block using
:class:`svg2ooxml.drawingml.animation.xml_builders.AnimationXMLBuilder` and
the timing-ID allocator.
"""

from __future__ import annotations

import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import Callable, Sequence

from lxml import etree as ET
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu, Inches

from svg2ooxml.drawingml.animation.id_allocator import TimingIDAllocator
from svg2ooxml.drawingml.animation.xml_builders import AnimationXMLBuilder
from svg2ooxml.drawingml.xml_builder import NS_P, to_string

SLIDE_WIDTH_IN = 10.0
SLIDE_HEIGHT_IN = 7.5

# Default hero shape positioned near slide center.
DEFAULT_SHAPE_LEFT_IN = 3.75
DEFAULT_SHAPE_TOP_IN = 2.75
DEFAULT_SHAPE_WIDTH_IN = 2.5
DEFAULT_SHAPE_HEIGHT_IN = 2.0


def _blank_layout(presentation) -> object:
    # Layout 6 is the plain blank layout in the default theme.
    return presentation.slide_layouts[6]


def new_presentation_with_hero_shape(
    *,
    fill: tuple[int, int, int] = (30, 110, 220),
    line: tuple[int, int, int] = (20, 20, 40),
    label: str | None = None,
) -> tuple[object, object, int]:
    """Create a 10x7.5in presentation with a blank slide 1 and a hero shape on slide 2.

    PowerPoint's build engine only pre-hides entrance-animation shapes on
    *slide transition*, not when a slideshow first starts. Hosting the
    animated shape on slide 2 and requiring a forward-advance to reach it
    guarantees the pre-hide state is applied before the click-group fires.

    Returns ``(presentation, animated_slide, shape_id)``. *shape_id* is the
    numeric ``spid`` later targeted by timing XML. The animated slide is
    ``presentation.slides[1]`` (index 0 is the blank intro slide).
    """
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)

    # Slide 1 — blank intro slide so slide 2 is entered via transition.
    presentation.slides.add_slide(_blank_layout(presentation))

    # Slide 2 — animated hero shape.
    slide = presentation.slides.add_slide(_blank_layout(presentation))
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(DEFAULT_SHAPE_LEFT_IN),
        Inches(DEFAULT_SHAPE_TOP_IN),
        Inches(DEFAULT_SHAPE_WIDTH_IN),
        Inches(DEFAULT_SHAPE_HEIGHT_IN),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(*line)

    if label:
        tf = shape.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        para.text = label
        para.runs[0].font.bold = True

    return presentation, slide, shape.shape_id


def build_timing_xml(
    par_factory: Callable[[AnimationXMLBuilder, int, int], ET._Element],
    *,
    animated_shape_ids: Sequence[str],
    start_id: int,
) -> str:
    """Allocate IDs, run a factory to build one ``<p:par>``, serialize timing tree."""
    xml_builder = AnimationXMLBuilder()
    allocator = TimingIDAllocator()
    ids = allocator.allocate(n_animations=1, start_id=start_id)
    anim_ids = ids.animations[0]
    par = par_factory(xml_builder, anim_ids.par, anim_ids.behavior)
    timing = xml_builder.build_timing_tree(
        ids=ids,
        animation_elements=[par],
        animated_shape_ids=list(animated_shape_ids),
    )
    return to_string(timing)


def inject_timing_into_pptx(
    pptx_path: Path,
    timing_xml: str,
    *,
    slide_filename: str = "ppt/slides/slide2.xml",
) -> None:
    """Replace *slide_filename*'s ``<p:timing>`` block with *timing_xml*."""
    temp_dir = Path(tempfile.mkdtemp(prefix="svg2ooxml-sample-"))
    temp_pptx = temp_dir / pptx_path.name
    try:
        with zipfile.ZipFile(pptx_path, "r") as source, zipfile.ZipFile(
            temp_pptx,
            "w",
            compression=zipfile.ZIP_DEFLATED,
        ) as target:
            for info in source.infolist():
                payload = source.read(info.filename)
                if info.filename == slide_filename:
                    payload = _patch_slide_xml(payload, timing_xml)
                target.writestr(info, payload)
        shutil.move(temp_pptx, pptx_path)
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def _patch_slide_xml(slide_xml: bytes, timing_xml: str) -> bytes:
    root = ET.fromstring(slide_xml)
    for existing in root.findall(f"{{{NS_P}}}timing"):
        root.remove(existing)
    timing_elem = ET.fromstring(_wrap_with_p_namespace(timing_xml).encode("utf-8"))
    root.append(timing_elem)
    return ET.tostring(root, encoding="utf-8", xml_declaration=True, standalone="yes")


def _wrap_with_p_namespace(timing_xml: str) -> str:
    """Inject the standard OOXML namespace declarations onto the root tag.

    ``to_string`` in ``drawingml.xml_builder`` strips the standard xmlns
    declarations because it assumes the emitted fragment will be grafted into
    a document that already declares them. For a free-standing parse we need
    to restore them so lxml can resolve both ``p:`` and ``a:`` prefixed tags.
    """
    prefix = "<p:timing"
    if prefix not in timing_xml:
        return timing_xml
    ns_decls = (
        ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    )
    return timing_xml.replace(prefix, prefix + ns_decls, 1)


__all__ = [
    "SLIDE_WIDTH_IN",
    "SLIDE_HEIGHT_IN",
    "new_presentation_with_hero_shape",
    "build_timing_xml",
    "inject_timing_into_pptx",
]
