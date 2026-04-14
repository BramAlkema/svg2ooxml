#!/usr/bin/env python3
"""Build a side-by-side W3C proof deck with browser and PowerPoint captures."""

from __future__ import annotations

import argparse
import json
import logging
import math
import shutil
from collections.abc import Sequence
from dataclasses import asdict, dataclass
from datetime import datetime
from pathlib import Path

from lxml import etree as ET
from PIL import Image, ImageDraw, ImageFont, ImageOps

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - dependency check
    Presentation = None  # type: ignore[assignment]
    RGBColor = None  # type: ignore[assignment]
    PP_ALIGN = None  # type: ignore[assignment]
    Inches = None  # type: ignore[assignment]
    Pt = None  # type: ignore[assignment]
    _PPTX_IMPORT_ERROR = exc
else:
    _PPTX_IMPORT_ERROR = None

from svg2ooxml.core.animation.parser import SMILParser
from svg2ooxml.core.animation.sampler import TimelineSampler
from tools.visual.browser_renderer import (
    BrowserRenderError,
    BrowserSvgRenderer,
    default_browser_renderer,
)
from tools.visual.builder import PptxBuilder
from tools.visual.renderer import PowerPointRenderer, VisualRendererError
from tools.visual.w3c_animation_suite import SCENARIOS as ANIMATION_SCENARIOS
from tools.visual.w3c_suite import SCENARIOS as STATIC_SCENARIOS

logger = logging.getLogger(__name__)

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
MAX_MONTAGE_FRAMES = 6
VARIANT_SPECS = {
    "browser": ("Browser SVG", None),
    "native": ("Native", "direct"),
    "mimic": ("Mimic", "mimic"),
    "rasterised": ("Rasterised", "bitmap"),
}
DEFAULT_VARIANTS = ("native",)
ALL_VARIANTS = ("browser", "native", "mimic", "rasterised")


@dataclass(frozen=True)
class ScenarioSpec:
    name: str
    svg_path: Path
    animated: bool


@dataclass
class VariantArtifacts:
    key: str
    label: str
    status: str
    preview_path: str
    pptx_path: str | None = None
    apng_path: str | None = None
    error: str | None = None


@dataclass
class ScenarioArtifacts:
    name: str
    svg_path: str
    animated: bool
    duration: float | None
    fps: float | None
    artifact_dir: str
    variants: list[VariantArtifacts]


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("reports/visual/w3c-proof-deck"),
        help="Output directory for the generated proof deck and artifacts.",
    )
    parser.add_argument(
        "--static-scenarios",
        nargs="*",
        default=None,
        help="Optional explicit static scenario names. Defaults to all static W3C scenarios.",
    )
    parser.add_argument(
        "--animation-scenarios",
        nargs="*",
        default=None,
        help="Optional explicit animation scenario names. Defaults to the animation suite.",
    )
    parser.add_argument(
        "--skip-static",
        action="store_true",
        help="Skip static scenarios entirely.",
    )
    parser.add_argument(
        "--skip-animations",
        action="store_true",
        help="Skip animation scenarios entirely.",
    )
    parser.add_argument(
        "--static-limit",
        type=int,
        default=None,
        help="Optional maximum number of static scenarios to include.",
    )
    parser.add_argument(
        "--animation-limit",
        type=int,
        default=None,
        help="Optional maximum number of animated scenarios to include.",
    )
    parser.add_argument(
        "--animation-duration",
        type=float,
        default=4.0,
        help="Default animation capture duration in seconds.",
    )
    parser.add_argument(
        "--max-animation-duration",
        type=float,
        default=5.0,
        help="Cap auto-detected animation durations to keep proof runs bounded.",
    )
    parser.add_argument(
        "--fps",
        type=float,
        default=4.0,
        help="Frame rate for live animation capture and APNG export.",
    )
    parser.add_argument(
        "--montage-frames",
        type=int,
        default=MAX_MONTAGE_FRAMES,
        help="Maximum frames to show in each animation montage.",
    )
    parser.add_argument(
        "--variants",
        nargs="+",
        choices=tuple(VARIANT_SPECS),
        default=list(DEFAULT_VARIANTS),
        help=(
            "Capture variants to include. Defaults to native only. "
            "Use --all-variants for browser/native/mimic/rasterised."
        ),
    )
    parser.add_argument(
        "--all-variants",
        action="store_true",
        help="Capture browser, native, mimic, and rasterised variants.",
    )
    parser.add_argument(
        "--deck-name",
        default="w3c-proof-deck.pptx",
        help="Filename for the generated proof deck.",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Enable debug logging.",
    )
    return parser.parse_args()


def _require_python_pptx() -> None:
    if Presentation is None or Inches is None or Pt is None or PP_ALIGN is None:
        raise RuntimeError(
            "python-pptx is required for proof deck generation."
        ) from _PPTX_IMPORT_ERROR


def _resolve_scenarios(
    all_scenarios: dict[str, Path],
    requested_names: Sequence[str] | None,
    *,
    limit: int | None,
    animated: bool,
) -> list[ScenarioSpec]:
    names = list(requested_names) if requested_names else sorted(all_scenarios)
    unknown = [name for name in names if name not in all_scenarios]
    if unknown:
        raise SystemExit(
            f"Unknown {'animation' if animated else 'static'} scenario(s): {', '.join(unknown)}"
        )
    if limit is not None:
        if limit < 0:
            raise SystemExit("Scenario limits must be >= 0.")
        names = names[:limit]
    return [
        ScenarioSpec(name=name, svg_path=all_scenarios[name], animated=animated)
        for name in names
    ]


def _detect_animation_duration(
    svg_text: str,
    *,
    default_duration: float,
    max_duration: float | None,
) -> float:
    try:
        parser = ET.XMLParser(recover=True)
        root = ET.fromstring(svg_text.encode("utf-8"), parser)
        if root is None:
            return default_duration
        animations = SMILParser().parse_svg_animations(root)
    except Exception:
        return default_duration

    if not animations:
        return default_duration

    try:
        summary = TimelineSampler().generate_keyframe_summary(animations)
        duration = float(summary.get("duration") or 0.0)
    except Exception:
        duration = 0.0

    if duration <= 0:
        duration = default_duration
    if max_duration is not None:
        duration = min(duration, max_duration)
    return max(0.1, duration)


def _select_sample_indices(total: int, max_items: int) -> list[int]:
    if total <= 0:
        return []
    if max_items <= 0:
        return []
    if total <= max_items:
        return list(range(total))
    if max_items == 1:
        return [0]
    last = total - 1
    chosen: list[int] = []
    for slot in range(max_items):
        index = round((slot * last) / (max_items - 1))
        if not chosen or index != chosen[-1]:
            chosen.append(index)
    if chosen[-1] != last:
        chosen[-1] = last
    return chosen


def _make_placeholder(
    output_path: Path,
    *,
    title: str,
    detail: str,
    size: tuple[int, int] = (1280, 720),
) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", size, color=(252, 252, 252))
    draw = ImageDraw.Draw(image)
    font = ImageFont.load_default()
    draw.rectangle((24, 24, size[0] - 24, size[1] - 24), outline=(180, 0, 0), width=4)
    draw.text((48, 48), title, font=font, fill=(150, 0, 0))
    wrapped = []
    line = ""
    for token in detail.split():
        candidate = token if not line else f"{line} {token}"
        if len(candidate) <= 88:
            line = candidate
            continue
        wrapped.append(line)
        line = token
    if line:
        wrapped.append(line)
    for row, chunk in enumerate(wrapped[:20], start=0):
        draw.text((48, 96 + (row * 18)), chunk, font=font, fill=(64, 64, 64))
    image.save(output_path)
    image.close()
    return output_path


def _copy_preview(source: Path, output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, output_path)
    return output_path


def _save_apng(frame_paths: Sequence[Path], output_path: Path, *, fps: float) -> Path:
    if not frame_paths:
        raise ValueError("Cannot create APNG without frames.")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    images = [Image.open(path).convert("RGBA") for path in frame_paths]
    try:
        duration_ms = max(1, int(round(1000.0 / max(fps, 0.001))))
        images[0].save(
            output_path,
            save_all=True,
            append_images=images[1:],
            duration=duration_ms,
            loop=0,
            disposal=2,
        )
    finally:
        for image in images:
            image.close()
    return output_path


def _build_montage(
    frame_paths: Sequence[Path],
    output_path: Path,
    *,
    fps: float,
    max_frames: int,
) -> Path:
    if not frame_paths:
        raise ValueError("Cannot create montage without frames.")

    selected_indices = _select_sample_indices(len(frame_paths), max_frames)
    selected_paths = [frame_paths[index] for index in selected_indices]
    images = [Image.open(path).convert("RGB") for path in selected_paths]
    try:
        max_width = max(image.width for image in images)
        max_height = max(image.height for image in images)
        thumb_box = (min(max_width, 320), min(max_height, 240))
        columns = 2 if len(images) > 1 else 1
        rows = int(math.ceil(len(images) / columns))
        padding = 12
        label_height = 20
        canvas_width = padding + columns * (thumb_box[0] + padding)
        canvas_height = padding + rows * (thumb_box[1] + label_height + padding)
        canvas = Image.new("RGB", (canvas_width, canvas_height), color=(255, 255, 255))
        draw = ImageDraw.Draw(canvas)
        font = ImageFont.load_default()

        for slot, (frame_index, image) in enumerate(zip(selected_indices, images, strict=False)):
            row = slot // columns
            column = slot % columns
            left = padding + column * (thumb_box[0] + padding)
            top = padding + row * (thumb_box[1] + label_height + padding)
            thumb = ImageOps.contain(image, thumb_box)
            thumb_left = left + (thumb_box[0] - thumb.width) // 2
            thumb_top = top
            canvas.paste(thumb, (thumb_left, thumb_top))
            draw.rectangle(
                (
                    left,
                    top,
                    left + thumb_box[0],
                    top + thumb_box[1],
                ),
                outline=(180, 180, 180),
                width=1,
            )
            timestamp = frame_index / max(fps, 0.001)
            draw.text(
                (left + 4, top + thumb_box[1] + 3),
                f"t={timestamp:.2f}s",
                font=font,
                fill=(70, 70, 70),
            )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        canvas.save(output_path)
        canvas.close()
        return output_path
    finally:
        for image in images:
            image.close()


def _fit_image_dimensions(
    image_path: Path,
    *,
    max_width_in: float,
    max_height_in: float,
) -> tuple[float, float]:
    with Image.open(image_path) as image:
        width_px = max(1, image.width)
        height_px = max(1, image.height)
    scale = min(max_width_in / width_px, max_height_in / height_px)
    return width_px * scale, height_px * scale


def _add_fitted_picture(
    slide,
    image_path: Path,
    *,
    left_in: float,
    top_in: float,
    box_width_in: float,
    box_height_in: float,
) -> None:
    width_in, height_in = _fit_image_dimensions(
        image_path,
        max_width_in=box_width_in,
        max_height_in=box_height_in,
    )
    offset_x = left_in + ((box_width_in - width_in) / 2.0)
    offset_y = top_in + ((box_height_in - height_in) / 2.0)
    slide.shapes.add_picture(
        str(image_path),
        Inches(offset_x),
        Inches(offset_y),
        width=Inches(width_in),
        height=Inches(height_in),
    )


def _relativize(path: Path | None, root: Path) -> str | None:
    if path is None:
        return None
    try:
        return str(path.relative_to(root))
    except ValueError:
        return str(path)


def _render_browser_variant(
    browser_renderer: BrowserSvgRenderer,
    scenario: ScenarioSpec,
    svg_text: str,
    *,
    duration: float | None,
    fps: float,
    output_dir: Path,
    montage_frames: int,
    root_output_dir: Path,
) -> VariantArtifacts:
    try:
        if scenario.animated:
            frames_dir = output_dir / "frames"
            frame_paths = browser_renderer.capture_animation(
                svg_text,
                frames_dir,
                duration=duration or 0.0,
                fps=fps,
                source_path=scenario.svg_path,
            )
            preview_path = _build_montage(
                frame_paths,
                output_dir / "preview.png",
                fps=fps,
                max_frames=montage_frames,
            )
            apng_path = _save_apng(frame_paths, output_dir / "preview.apng", fps=fps)
        else:
            rendered = browser_renderer.render_svg(
                svg_text,
                output_dir / "preview.png",
                source_path=scenario.svg_path,
            )
            preview_path = rendered.image
            apng_path = None
        return VariantArtifacts(
            key="browser",
            label="Browser SVG",
            status="ok",
            preview_path=_relativize(preview_path, root_output_dir) or str(preview_path),
            apng_path=_relativize(apng_path, root_output_dir),
        )
    except (BrowserRenderError, OSError, RuntimeError, ValueError) as exc:
        placeholder = _make_placeholder(
            output_dir / "preview.png",
            title="Browser capture failed",
            detail=str(exc),
        )
        return VariantArtifacts(
            key="browser",
            label="Browser SVG",
            status="error",
            preview_path=_relativize(placeholder, root_output_dir) or str(placeholder),
            error=str(exc),
        )


def _render_powerpoint_variant(
    renderer: PowerPointRenderer,
    scenario: ScenarioSpec,
    svg_text: str,
    *,
    fidelity_tier: str,
    duration: float | None,
    fps: float,
    output_dir: Path,
    montage_frames: int,
    root_output_dir: Path,
) -> VariantArtifacts:
    pptx_path = output_dir / "presentation.pptx"
    try:
        builder = PptxBuilder(slide_size_mode="same", fidelity_tier=fidelity_tier)
        builder.build_from_svg(svg_text, pptx_path, source_path=scenario.svg_path)
        if scenario.animated:
            frame_paths = renderer.capture_animation(
                pptx_path,
                output_dir / "frames",
                duration=duration or 0.0,
                fps=fps,
            )
            preview_path = _build_montage(
                frame_paths,
                output_dir / "preview.png",
                fps=fps,
                max_frames=montage_frames,
            )
            apng_path = _save_apng(frame_paths, output_dir / "preview.apng", fps=fps)
        else:
            rendered = renderer.render(pptx_path, output_dir / "render")
            preview_path = _copy_preview(rendered.images[0], output_dir / "preview.png")
            apng_path = None
        return VariantArtifacts(
            key=output_dir.name,
            label=output_dir.name.replace("_", " ").title(),
            status="ok",
            preview_path=_relativize(preview_path, root_output_dir) or str(preview_path),
            pptx_path=_relativize(pptx_path, root_output_dir),
            apng_path=_relativize(apng_path, root_output_dir),
        )
    except (VisualRendererError, OSError, RuntimeError, ValueError) as exc:
        placeholder = _make_placeholder(
            output_dir / "preview.png",
            title=f"{output_dir.name} capture failed",
            detail=str(exc),
        )
        return VariantArtifacts(
            key=output_dir.name,
            label=output_dir.name.replace("_", " ").title(),
            status="error",
            preview_path=_relativize(placeholder, root_output_dir) or str(placeholder),
            pptx_path=_relativize(pptx_path, root_output_dir) if pptx_path.exists() else None,
            error=str(exc),
        )


def _build_scenario_artifacts(
    scenario: ScenarioSpec,
    *,
    output_dir: Path,
    browser_renderer: BrowserSvgRenderer | None,
    powerpoint_renderer: PowerPointRenderer,
    variant_keys: Sequence[str],
    default_animation_duration: float,
    max_animation_duration: float | None,
    fps: float,
    montage_frames: int,
) -> ScenarioArtifacts:
    scenario_dir = output_dir / "artifacts" / scenario.name
    scenario_dir.mkdir(parents=True, exist_ok=True)
    svg_text = scenario.svg_path.read_text(encoding="utf-8")
    duration = None
    if scenario.animated:
        duration = _detect_animation_duration(
            svg_text,
            default_duration=default_animation_duration,
            max_duration=max_animation_duration,
        )

    variants: list[VariantArtifacts] = []
    if "browser" in variant_keys:
        if browser_renderer is None:
            raise RuntimeError("Browser variant requested but browser renderer is unavailable.")
        variants.append(
            _render_browser_variant(
                browser_renderer,
                scenario,
                svg_text,
                duration=duration,
                fps=fps,
                output_dir=scenario_dir / "browser",
                montage_frames=montage_frames,
                root_output_dir=output_dir,
            )
        )

    for key in variant_keys:
        if key == "browser":
            continue
        label, fidelity_tier = VARIANT_SPECS[key]
        output_key = key.replace(" ", "_")
        variant = _render_powerpoint_variant(
            powerpoint_renderer,
            scenario,
            svg_text,
            fidelity_tier=fidelity_tier or "direct",
            duration=duration,
            fps=fps,
            output_dir=scenario_dir / output_key,
            montage_frames=montage_frames,
            root_output_dir=output_dir,
        )
        variant.key = key
        variant.label = label
        variants.append(variant)

    return ScenarioArtifacts(
        name=scenario.name,
        svg_path=_relativize(scenario.svg_path, output_dir) or str(scenario.svg_path),
        animated=scenario.animated,
        duration=duration,
        fps=fps if scenario.animated else None,
        artifact_dir=_relativize(scenario_dir, output_dir) or str(scenario_dir),
        variants=variants,
    )


def _write_manifest(results: Sequence[ScenarioArtifacts], output_dir: Path) -> Path:
    manifest_path = output_dir / "manifest.json"
    payload = {
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "scenarios": [asdict(item) for item in results],
    }
    manifest_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    return manifest_path


def _add_textbox(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    bold: bool = False,
    color: tuple[int, int, int] = (0, 0, 0),
    align: str | None = None,
) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    if align == "center":
        paragraph.alignment = PP_ALIGN.CENTER
    elif align == "right":
        paragraph.alignment = PP_ALIGN.RIGHT
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def _build_proof_deck(
    results: Sequence[ScenarioArtifacts],
    *,
    output_dir: Path,
    deck_name: str,
) -> Path:
    _require_python_pptx()
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = presentation.slide_layouts[6]

    title_slide = presentation.slides.add_slide(blank_layout)
    _add_textbox(
        title_slide,
        left=0.5,
        top=0.5,
        width=12.3,
        height=0.6,
        text="W3C SVG Conversion Proof Deck",
        font_size=28,
        bold=True,
    )
    _add_textbox(
        title_slide,
        left=0.5,
        top=1.2,
        width=12.3,
        height=0.8,
        text=(
            "Columns: "
            + ", ".join(
                variant.label for result in results[:1] for variant in result.variants
            )
            + ". "
            "Animated slides show frame montages; APNG artifacts live under the run output."
        ),
        font_size=16,
    )
    _add_textbox(
        title_slide,
        left=0.5,
        top=2.1,
        width=12.3,
        height=0.4,
        text=f"Generated {datetime.now().isoformat(timespec='seconds')}",
        font_size=14,
        color=(90, 90, 90),
    )

    margin_left = 0.25
    margin_right = 0.25
    gap = 0.12

    for result in results:
        variant_count = max(1, len(result.variants))
        column_width = (
            SLIDE_WIDTH_IN
            - margin_left
            - margin_right
            - (gap * (variant_count - 1))
        ) / variant_count
        slide = presentation.slides.add_slide(blank_layout)
        _add_textbox(
            slide,
            left=0.25,
            top=0.18,
            width=10.8,
            height=0.45,
            text=result.name,
            font_size=22,
            bold=True,
        )
        source_line = f"{result.svg_path}"
        if result.animated and result.duration is not None and result.fps is not None:
            source_line += f" | animated capture {result.duration:.2f}s @ {result.fps:.1f} fps"
        _add_textbox(
            slide,
            left=0.25,
            top=0.62,
            width=12.8,
            height=0.32,
            text=source_line,
            font_size=11,
            color=(90, 90, 90),
        )

        for index, variant in enumerate(result.variants):
            left = margin_left + index * (column_width + gap)
            _add_textbox(
                slide,
                left=left,
                top=0.95,
                width=column_width,
                height=0.3,
                text=variant.label,
                font_size=15,
                bold=True,
                align="center",
            )
            preview_path = output_dir / variant.preview_path
            _add_fitted_picture(
                slide,
                preview_path,
                left_in=left,
                top_in=1.25,
                box_width_in=column_width,
                box_height_in=5.2,
            )
            status = variant.status.upper()
            footer_bits = [status]
            if variant.apng_path:
                footer_bits.append(Path(variant.apng_path).name)
            _add_textbox(
                slide,
                left=left,
                top=6.55,
                width=column_width,
                height=0.42,
                text=" | ".join(footer_bits),
                font_size=10,
                align="center",
                color=(90, 90, 90) if variant.status == "ok" else (150, 0, 0),
            )

        _add_textbox(
            slide,
            left=0.25,
            top=7.0,
            width=12.8,
            height=0.26,
            text=f"Artifacts: {result.artifact_dir}",
            font_size=10,
            color=(90, 90, 90),
        )

    deck_path = output_dir / deck_name
    output_dir.mkdir(parents=True, exist_ok=True)
    presentation.save(deck_path)
    return deck_path


def _iter_scenarios(args: argparse.Namespace) -> list[ScenarioSpec]:
    scenarios: list[ScenarioSpec] = []
    if not args.skip_static:
        scenarios.extend(
            _resolve_scenarios(
                STATIC_SCENARIOS,
                args.static_scenarios,
                limit=args.static_limit,
                animated=False,
            )
        )
    if not args.skip_animations:
        scenarios.extend(
            _resolve_scenarios(
                ANIMATION_SCENARIOS,
                args.animation_scenarios,
                limit=args.animation_limit,
                animated=True,
            )
        )
    if not scenarios:
        raise SystemExit("No scenarios selected.")
    return scenarios


def _selected_variant_keys(args: argparse.Namespace) -> tuple[str, ...]:
    raw_keys = ALL_VARIANTS if args.all_variants else tuple(args.variants)
    selected: list[str] = []
    for key in raw_keys:
        if key in selected:
            continue
        selected.append(key)
    if not selected:
        raise SystemExit("At least one capture variant must be selected.")
    return tuple(selected)


def main() -> None:
    args = _parse_args()
    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(levelname)s %(message)s",
    )
    _require_python_pptx()

    output_dir = args.output.resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    variant_keys = _selected_variant_keys(args)

    browser_renderer = default_browser_renderer() if "browser" in variant_keys else None
    if browser_renderer is not None and not browser_renderer.available:
        raise SystemExit(
            "Browser rendering is unavailable. Install Playwright and its browser runtime."
        )

    powerpoint_renderer = PowerPointRenderer()
    if not powerpoint_renderer.available:
        raise SystemExit(
            "PowerPoint slideshow capture is unavailable. This tool requires macOS PowerPoint automation."
        )

    scenarios = _iter_scenarios(args)
    logger.info("Selected %d scenarios.", len(scenarios))

    results: list[ScenarioArtifacts] = []
    for index, scenario in enumerate(scenarios, start=1):
        logger.info(
            "[%d/%d] %s%s",
            index,
            len(scenarios),
            scenario.name,
            " (animated)" if scenario.animated else "",
        )
        results.append(
            _build_scenario_artifacts(
                scenario,
                output_dir=output_dir,
                browser_renderer=browser_renderer,
                powerpoint_renderer=powerpoint_renderer,
                variant_keys=variant_keys,
                default_animation_duration=args.animation_duration,
                max_animation_duration=args.max_animation_duration,
                fps=args.fps,
                montage_frames=args.montage_frames,
            )
        )

    manifest_path = _write_manifest(results, output_dir)
    deck_path = _build_proof_deck(results, output_dir=output_dir, deck_name=args.deck_name)
    logger.info("Wrote manifest: %s", manifest_path)
    logger.info("Wrote proof deck: %s", deck_path)


if __name__ == "__main__":
    main()
