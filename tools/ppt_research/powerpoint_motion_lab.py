#!/usr/bin/env python3
"""Build a native PowerPoint motion lab deck for slideshow semantics checks."""

from __future__ import annotations

import argparse
import logging
import shutil
import tempfile
import zipfile
from collections.abc import Iterator, Sequence
from dataclasses import dataclass
from itertools import count
from pathlib import Path

from lxml import etree as ET

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - dependency check
    Presentation = None  # type: ignore[assignment]
    RGBColor = None  # type: ignore[assignment]
    MSO_AUTO_SHAPE_TYPE = None  # type: ignore[assignment]
    MSO_CONNECTOR = None  # type: ignore[assignment]
    PP_ALIGN = None  # type: ignore[assignment]
    Inches = None  # type: ignore[assignment]
    Pt = None  # type: ignore[assignment]
    _PPTX_IMPORT_ERROR = exc
else:
    _PPTX_IMPORT_ERROR = None

from svg2ooxml.drawingml.animation.id_allocator import TimingIDAllocator
from svg2ooxml.drawingml.animation.xml_builders import AnimationXMLBuilder
from svg2ooxml.drawingml.writer import DEFAULT_SLIDE_SIZE
from svg2ooxml.drawingml.xml_builder import NS_P, p_elem

logger = logging.getLogger(__name__)

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5


@dataclass(frozen=True, slots=True)
class MotionLabCase:
    key: str
    label: str
    motion_dx_in: float = 0.0
    motion_dy_in: float = 0.0
    initial_rotation_deg: float | None = None
    motion_duration_ms: int = 2200
    initial_rotation_ms: int = 1
    child_order: str = "motion-first"

    @property
    def has_motion(self) -> bool:
        return abs(self.motion_dx_in) > 1e-9 or abs(self.motion_dy_in) > 1e-9

    @property
    def has_rotation(self) -> bool:
        return self.initial_rotation_deg is not None and abs(self.initial_rotation_deg) > 1e-9


@dataclass(frozen=True, slots=True)
class PlacedMotionLabCase:
    case: MotionLabCase
    triangle_shape_id: int


def _require_python_pptx() -> None:
    if Presentation is None or Inches is None or Pt is None:
        raise RuntimeError(
            "python-pptx is required to build the PowerPoint motion lab."
        ) from _PPTX_IMPORT_ERROR


def _case_specs() -> list[MotionLabCase]:
    return [
        MotionLabCase("rot_neg90", "Rot only: -90", initial_rotation_deg=-90.0),
        MotionLabCase("rot_pos90", "Rot only: +90", initial_rotation_deg=90.0),
        MotionLabCase("down_only", "Motion only: down", motion_dy_in=0.95),
        MotionLabCase(
            "down_rot_neg90",
            "Down + rot -90",
            motion_dy_in=0.95,
            initial_rotation_deg=-90.0,
        ),
        MotionLabCase(
            "down_rot_pos90",
            "Down + rot +90",
            motion_dy_in=0.95,
            initial_rotation_deg=90.0,
        ),
        MotionLabCase(
            "right_rot_neg90",
            "Right + rot -90",
            motion_dx_in=1.2,
            initial_rotation_deg=-90.0,
        ),
        MotionLabCase(
            "right_rot_pos90_motion_first",
            "Right + rot +90 (M then R)",
            motion_dx_in=1.2,
            initial_rotation_deg=90.0,
            child_order="motion-first",
        ),
        MotionLabCase(
            "right_rot_pos90_rot_first",
            "Right + rot +90 (R then M)",
            motion_dx_in=1.2,
            initial_rotation_deg=90.0,
            child_order="rotation-first",
        ),
    ]


def _format_fraction(value: float) -> str:
    if abs(value) < 1e-10:
        return "0"
    return f"{value:.6g}"


def _relative_motion_path(
    *,
    dx_in: float,
    dy_in: float,
    slide_width_emu: int,
    slide_height_emu: int,
) -> str:
    dx_frac = Inches(dx_in).emu / slide_width_emu
    dy_frac = Inches(dy_in).emu / slide_height_emu
    return f"M 0 0 L {_format_fraction(dx_frac)} {_format_fraction(dy_frac)} E"


def _build_rotation_child(
    *,
    xml_builder: AnimationXMLBuilder,
    shape_id: int,
    behavior_id: int,
    rotation_deg: float,
    duration_ms: int,
) -> ET._Element:
    anim_rot = p_elem(
        "animRot",
        by=str(int(round(rotation_deg * 60000))),
    )
    anim_rot.append(
        xml_builder.build_behavior_core_elem(
            behavior_id=behavior_id,
            duration_ms=max(1, duration_ms),
            target_shape=str(shape_id),
            attr_name_list=["r"],
            fill_mode="freeze",
        )
    )
    return anim_rot


def _build_motion_child(
    *,
    xml_builder: AnimationXMLBuilder,
    shape_id: int,
    behavior_id: int,
    motion_path: str,
    duration_ms: int,
) -> ET._Element:
    anim_motion = p_elem(
        "animMotion",
        origin="layout",
        path=motion_path,
        pathEditMode="relative",
    )
    anim_motion.append(
        xml_builder.build_behavior_core_elem(
            behavior_id=behavior_id,
            duration_ms=max(1, duration_ms),
            target_shape=str(shape_id),
            fill_mode="freeze",
        )
    )
    return anim_motion


def _build_case_par(
    *,
    xml_builder: AnimationXMLBuilder,
    placed_case: PlacedMotionLabCase,
    par_id: int,
    behavior_id: int,
    id_counter: Iterator[int],
    slide_size: tuple[int, int],
) -> ET._Element:
    case = placed_case.case
    shape_id = placed_case.triangle_shape_id

    if case.has_motion:
        motion_child = _build_motion_child(
            xml_builder=xml_builder,
            shape_id=shape_id,
            behavior_id=behavior_id,
            motion_path=_relative_motion_path(
                dx_in=case.motion_dx_in,
                dy_in=case.motion_dy_in,
                slide_width_emu=slide_size[0],
                slide_height_emu=slide_size[1],
            ),
            duration_ms=case.motion_duration_ms,
        )
        child_elements: list[ET._Element] = [motion_child]
        if case.has_rotation:
            rotation_child = _build_rotation_child(
                xml_builder=xml_builder,
                shape_id=shape_id,
                behavior_id=next(id_counter),
                rotation_deg=case.initial_rotation_deg or 0.0,
                duration_ms=case.initial_rotation_ms,
            )
            rotation_par = xml_builder.build_delayed_child_par(
                par_id=next(id_counter),
                delay_ms=0,
                duration_ms=case.initial_rotation_ms,
                child_element=rotation_child,
            )
            if case.child_order == "rotation-first":
                child_elements = [rotation_par, motion_child]
            else:
                child_elements.append(rotation_par)
        return xml_builder.build_par_container_with_children_elem(
            par_id=par_id,
            duration_ms=case.motion_duration_ms,
            delay_ms=0,
            child_elements=child_elements,
            preset_class="path",
            node_type="clickEffect",
            default_target_shape=str(shape_id),
            effect_group_id=par_id,
        )

    rotation_child = _build_rotation_child(
        xml_builder=xml_builder,
        shape_id=shape_id,
        behavior_id=behavior_id,
        rotation_deg=case.initial_rotation_deg or 0.0,
        duration_ms=case.initial_rotation_ms,
    )
    return xml_builder.build_par_container_elem(
        par_id=par_id,
        duration_ms=case.initial_rotation_ms,
        delay_ms=0,
        child_element=rotation_child,
        preset_id=8,
        preset_class="emph",
        preset_subtype=0,
        node_type="clickEffect",
        default_target_shape=str(shape_id),
        effect_group_id=par_id,
    )


def _build_lab_timing(
    placed_cases: Sequence[PlacedMotionLabCase],
    *,
    start_id: int,
    slide_size: tuple[int, int],
) -> str:
    xml_builder = AnimationXMLBuilder()
    allocator = TimingIDAllocator()
    ids = allocator.allocate(len(placed_cases), start_id=start_id)
    next_available_id = ids.animations[-1].behavior + 1 if ids.animations else start_id + 3
    extra_ids = count(next_available_id)

    par_elements = [
        _build_case_par(
            xml_builder=xml_builder,
            placed_case=placed_case,
            par_id=anim_ids.par,
            behavior_id=anim_ids.behavior,
            id_counter=extra_ids,
            slide_size=slide_size,
        )
        for placed_case, anim_ids in zip(placed_cases, ids.animations, strict=True)
    ]
    timing = xml_builder.build_timing_tree(
        ids=ids,
        animation_elements=par_elements,
        animated_shape_ids=[str(case.triangle_shape_id) for case in placed_cases],
    )
    return ET.tostring(timing, encoding="unicode")


def _patch_slide_xml_with_timing(slide_xml: bytes, timing_xml: str) -> bytes:
    slide_root = ET.fromstring(slide_xml)
    for existing in slide_root.findall(f"{{{NS_P}}}timing"):
        slide_root.remove(existing)
    timing_elem = ET.fromstring(timing_xml.encode("utf-8"))
    slide_root.append(timing_elem)
    return ET.tostring(slide_root, encoding="utf-8", xml_declaration=True, standalone="yes")


def _inject_timing_into_pptx(pptx_path: Path, timing_xml: str) -> None:
    temp_dir = Path(tempfile.mkdtemp(prefix="ppt-motion-lab-"))
    temp_pptx = temp_dir / pptx_path.name
    try:
        with zipfile.ZipFile(pptx_path, "r") as source, zipfile.ZipFile(
            temp_pptx,
            "w",
            compression=zipfile.ZIP_DEFLATED,
        ) as target:
            for info in source.infolist():
                payload = source.read(info.filename)
                if info.filename == "ppt/slides/slide1.xml":
                    payload = _patch_slide_xml_with_timing(payload, timing_xml)
                target.writestr(info, payload)
        shutil.move(temp_pptx, pptx_path)
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def _add_textbox(slide, left, top, width, height, text: str, *, font_size: int, bold: bool = False, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    return box


def _add_case_shapes(slide, case: MotionLabCase, *, col: int, row: int) -> PlacedMotionLabCase:
    cell_left = 0.45 + col * 3.15
    cell_top = 1.15 + row * 2.7
    start_x = cell_left + 0.35
    start_y = cell_top + 0.65
    end_x = start_x + max(case.motion_dx_in, 0.0)
    end_y = start_y + max(case.motion_dy_in, 0.0)

    if case.has_motion and case.motion_dx_in < 0:
        end_x = start_x + case.motion_dx_in
    if case.has_motion and case.motion_dy_in < 0:
        end_y = start_y + case.motion_dy_in

    _add_textbox(
        slide,
        Inches(cell_left),
        Inches(cell_top),
        Inches(2.8),
        Inches(0.35),
        case.label,
        font_size=13,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    start_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(start_x),
        Inches(start_y + 0.42),
        Inches(0.34),
        Inches(0.34),
    )
    start_box.fill.solid()
    start_box.fill.fore_color.rgb = RGBColor(255, 209, 209)
    start_box.line.color.rgb = RGBColor(0, 0, 0)
    start_box.line.width = Pt(1.5)

    if case.has_motion:
        end_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(end_x),
            Inches(end_y + 0.42),
            Inches(0.34),
            Inches(0.34),
        )
        end_box.fill.solid()
        end_box.fill.fore_color.rgb = RGBColor(255, 209, 209)
        end_box.line.color.rgb = RGBColor(0, 0, 0)
        end_box.line.width = Pt(1.5)
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(start_x + 0.17),
            Inches(start_y + 0.59),
            Inches(end_x + 0.17),
            Inches(end_y + 0.59),
        )
        connector.line.color.rgb = RGBColor(80, 80, 80)
        connector.line.width = Pt(1.25)

    triangle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        Inches(start_x + 0.03),
        Inches(start_y),
        Inches(0.34),
        Inches(0.34),
    )
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor(16, 64, 255)
    triangle.line.color.rgb = RGBColor(24, 136, 24)
    triangle.line.width = Pt(1.75)

    return PlacedMotionLabCase(case=case, triangle_shape_id=triangle.shape_id)


def build_motion_lab_pptx(output_path: Path) -> Path:
    _require_python_pptx()

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    _add_textbox(
        slide,
        Inches(0.4),
        Inches(0.15),
        Inches(12.5),
        Inches(0.4),
        "PowerPoint Motion Lab",
        font_size=24,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.55),
        Inches(12.2),
        Inches(0.35),
        "Frame 0: check static sign and initial pose. Later frames: compare M→R vs R→M ordering.",
        font_size=12,
        align=PP_ALIGN.CENTER,
    )

    placed_cases: list[PlacedMotionLabCase] = []
    for index, case in enumerate(_case_specs()):
        placed_cases.append(
            _add_case_shapes(slide, case, col=index % 4, row=index // 4)
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)

    max_shape_id = max(case.triangle_shape_id for case in placed_cases)
    timing_xml = _build_lab_timing(
        placed_cases,
        start_id=max_shape_id + 1,
        slide_size=DEFAULT_SLIDE_SIZE,
    )
    _inject_timing_into_pptx(output_path, timing_xml)
    return output_path


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=Path("reports/visual/powerpoint/motion-lab/motion-lab.pptx"),
        help="Output PPTX path.",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Enable debug logging.",
    )
    return parser.parse_args()


def main() -> None:
    args = _parse_args()
    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(levelname)s %(message)s",
    )
    output = build_motion_lab_pptx(args.output)
    logger.info("Built motion lab deck: %s", output)


if __name__ == "__main__":
    main()
